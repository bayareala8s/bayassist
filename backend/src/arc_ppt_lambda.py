import os
import boto3
from io import BytesIO
from pptx import Presentation
from common.models import now_iso

dynamodb = boto3.resource("dynamodb")
s3 = boto3.client("s3")

JOBS_TABLE = os.environ["JOBS_TABLE"]
OUTPUT_BUCKET = os.environ["OUTPUT_BUCKET"]

jobs_table = dynamodb.Table(JOBS_TABLE)


def build_arc_presentation(topology: dict, job_id: str) -> bytes:
  prs = Presentation()

  slide_layout = prs.slide_layouts[0]
  slide = prs.slides.add_slide(slide_layout)
  slide.shapes.title.text = "BayAssist – Architecture Review"
  subtitle = slide.placeholders[1]
  subtitle.text = f"Job ID: {job_id}"

  layout = prs.slide_layouts[1]
  slide = prs.slides.add_slide(layout)
  slide.shapes.title.text = "Architecture Overview"
  body = slide.placeholders[1].text_frame
  body.text = "High-level AWS components:"
  for svc in topology.get("services", []):
    p = body.add_paragraph()
    p.text = f"- {svc.get('type')} – {svc.get('name')}"
    p.level = 1

  slide = prs.slides.add_slide(layout)
  slide.shapes.title.text = "RTO / RPO"
  tf = slide.placeholders[1].text_frame
  tf.text = "Sample targets (customize per system):"
  for row in [
    "Critical APIs: RTO 15m, RPO 5m",
    "File-transfer flows: RTO 30m, RPO 15m",
    "Reporting/analytics: RTO 4h, RPO 1h",
  ]:
    p = tf.add_paragraph()
    p.text = f"- {row}"
    p.level = 1

  slide = prs.slides.add_slide(layout)
  slide.shapes.title.text = "Security & Risk"
  tf = slide.placeholders[1].text_frame
  tf.text = "Key controls and risks:"
  for point in [
    "All data in transit via TLS 1.2+",
    "KMS-encrypted S3 buckets, least-privilege IAM",
    "GuardDuty / CloudTrail enabled",
    "Risks: Misconfig, missing monitoring, human error",
  ]:
    p = tf.add_paragraph()
    p.text = f"- {point}"
    p.level = 1

  bio = BytesIO()
  prs.save(bio)
  return bio.getvalue()


def lambda_handler(event, context):
  job_id = event["job_id"]
  topology = event["topology"]
  ppt_bytes = build_arc_presentation(topology, job_id)

  ppt_key = f"outputs/{job_id}/arc-review.pptx"
  s3.put_object(
    Bucket=OUTPUT_BUCKET,
    Key=ppt_key,
    Body=ppt_bytes,
    ContentType="application/vnd.openxmlformats-officedocument.presentationml.presentation",
  )

  jobs_table.update_item(
    Key={"job_id": job_id},
    UpdateExpression="SET updated_at = :u, output_ppt_s3_key = :p",
    ExpressionAttributeValues={
      ":u": now_iso(),
      ":p": ppt_key,
    },
  )

  return {"job_id": job_id, "ppt_s3_key": ppt_key}
